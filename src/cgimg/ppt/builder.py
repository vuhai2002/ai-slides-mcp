"""Assemble image files into a full-bleed PPTX. Image aspect should match the
slide aspect (16:9 -> 1672x941 images) for edge-to-edge with no bars/crop."""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches

# Slide dimensions (inches) per aspect.
_SLIDE_DIMS = {
    "16:9": (13.333, 7.5),
    "1:1": (7.5, 7.5),
    "3:4": (7.5, 10.0),
    "9:16": (7.5, 13.333),
}


def build_pptx(image_paths: list[str], out_path: str, aspect: str = "16:9") -> str:
    if not image_paths:
        raise ValueError("no images provided")
    dims = _SLIDE_DIMS.get(aspect.strip().lower())
    if dims is None:
        raise ValueError(f"unsupported ppt aspect {aspect!r}; use {list(_SLIDE_DIMS)}")
    w_in, h_in = dims

    prs = Presentation()
    prs.slide_width = Inches(w_in)
    prs.slide_height = Inches(h_in)
    blank = prs.slide_layouts[6]

    for img in image_paths:
        if not os.path.exists(img):
            raise FileNotFoundError(img)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img, 0, 0, prs.slide_width, prs.slide_height)

    out_dir = os.path.dirname(os.path.abspath(out_path))
    os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    return out_path
